from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _safe_text(value: Optional[Any]) -> str:
    if value is None:
        return "N/A"
    text = str(value).strip()
    return text if text else "N/A"


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullet_slide(prs: Presentation, title: str, bullets: List[str], subtitle: Optional[str] = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    body = slide.placeholders[1].text_frame
    body.clear()
    for idx, bullet in enumerate(bullets):
        if idx == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.LEFT


def _add_table_slide(prs: Presentation, title: str, rows: List[List[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if len(rows) < 2:
        rows = rows + [["No data available", ""]]
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(0.7), Inches(1.6), Inches(12.0), Inches(4.8))
    table = table_shape.table
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(14)
    table.rows[0].height = Inches(0.4)


def _add_distribution_slide(prs: Presentation, summary: Dict[str, Any]) -> None:
    distribution = summary.get("health_distribution", {})
    rows = [["Status", "Count"], ["Green", str(distribution.get("Green", 0))], ["Amber", str(distribution.get("Amber", 0))], ["Red", str(distribution.get("Red", 0))]]
    _add_table_slide(prs, "Health Distribution", rows)


def _add_escalation_slide(prs: Presentation, summary: Dict[str, Any]) -> None:
    projects = summary.get("projects_requiring_escalation", [])
    if projects:
        rows = [["Project", "Client", "Manager", "RAG", "Confidence"]]
        for item in projects:
            rows.append([_safe_text(item.get("project_name")), _safe_text(item.get("client")), _safe_text(item.get("project_manager")), _safe_text(item.get("overall_rag")), f"{item.get('confidence', 0.0):.2f}"])
    else:
        rows = [["Project", "Client", "Manager", "RAG", "Confidence"], ["No projects currently require escalation", "", "", "", ""]]
    _add_table_slide(prs, "Projects Requiring Escalation", rows)


def _add_risk_themes_slide(prs: Presentation, summary: Dict[str, Any]) -> None:
    themes = summary.get("cross_project_risk_themes", []) or ["No recurring risk themes detected."]
    _add_bullet_slide(prs, "Cross-Project Risk Themes", themes)


def _add_recommendations_slide(prs: Presentation, summary: Dict[str, Any]) -> None:
    recommendations = summary.get("executive_recommendations", []) or ["No recommendations available."]
    _add_bullet_slide(prs, "Executive Recommendations", recommendations)


def _add_methodology_slide(prs: Presentation, summary: Dict[str, Any]) -> None:
    bullets = [
        "Deterministic scoring based on parsed project data and RAG analysis outputs.",
        "Schedule, milestone, blocker, sentiment, and budget signals are evaluated independently.",
        "The deck uses the JSON export from monthly_synthesis.py as its source of truth.",
        "Budget status is reported as Not Available when no budget data is present.",
    ]
    _add_bullet_slide(prs, "Methodology", bullets)


def build_presentation(summary: Dict[str, Any], output_path: Path) -> Path:
    prs = Presentation()
    title = summary.get("portfolio_name", "Executive Project Health Report")
    _add_title_slide(prs, "Executive Project Health Report", f"Portfolio view for {title}")
    _add_bullet_slide(prs, "Portfolio Overview", [
        f"Portfolio Health: {_safe_text(summary.get('portfolio_health'))}",
        f"Projects Reviewed: {len(summary.get('projects', []))}",
        f"Health Distribution: Green {_safe_text(summary.get('health_distribution', {}).get('Green'))}, Amber {_safe_text(summary.get('health_distribution', {}).get('Amber'))}, Red {_safe_text(summary.get('health_distribution', {}).get('Red'))}",
    ])
    _add_distribution_slide(prs, summary)
    _add_escalation_slide(prs, summary)
    _add_risk_themes_slide(prs, summary)
    _add_recommendations_slide(prs, summary)
    _add_methodology_slide(prs, summary)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def generate_presentation_from_json(summary_path: Path, output_path: Optional[Path] = None) -> Path:
    summary_path = Path(summary_path)
    with summary_path.open("r", encoding="utf-8") as handle:
        summary = json.load(handle)
    if output_path is None:
        output_path = summary_path.with_name("Executive_Project_Health_Report.pptx")
    return build_presentation(summary, output_path)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate a PowerPoint executive report from the monthly synthesis JSON")
    parser.add_argument("summary_json", help="Path to portfolio_summary.json")
    parser.add_argument("--output", default=None, help="Destination PowerPoint path")
    args = parser.parse_args()

    output_path = generate_presentation_from_json(Path(args.summary_json), output_path=Path(args.output) if args.output else None)
    print(output_path)
